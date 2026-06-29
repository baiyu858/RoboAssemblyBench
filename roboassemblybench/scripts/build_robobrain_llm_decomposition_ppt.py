#!/usr/bin/env python3
"""Build a one-slide RoboBrain LLM task-decomposition presentation."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
ASSET_DIR = ROOT / "roboassemblybench" / "outputs" / "presentation_assets"
OUTPUT_DIR = ROOT / "roboassemblybench" / "outputs" / "presentations"
PPTX_PATH = OUTPUT_DIR / "robobrain_llm_task_decomposition_and_next_month_plan.pptx"
PREVIEW_PATH = OUTPUT_DIR / "robobrain_llm_task_decomposition_one_page.png"
PLAN_PREVIEW_PATH = OUTPUT_DIR / "robobrain_llm_next_month_plan.png"
UI_FRAME = ASSET_DIR / "frame_035.png"
SIM_FRAME = ASSET_DIR / "frame_175.png"

SLIDE_W = 13.333
SLIDE_H = 7.5
PX_W = 2560
PX_H = 1440

COLORS = {
    "bg": "F6F7F9",
    "surface": "FFFFFF",
    "ink": "172033",
    "muted": "687386",
    "line": "DDE2E8",
    "teal": "157A6E",
    "teal_light": "EAF7F3",
    "indigo": "5267D8",
    "indigo_light": "EEF1FF",
    "green": "2E8B65",
    "green_light": "EAF8F0",
    "amber": "C77416",
    "amber_light": "FFF4DF",
    "red": "C34747",
}

FONT = "Noto Sans CJK SC"
FONT_PATH = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD_PATH = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: str = "ink",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS.get(color, color))
    return box


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = "surface",
    line: str = "line",
    radius: bool = True,
    width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS.get(fill, fill))
    shape.line.color.rgb = rgb(COLORS.get(line, line))
    shape.line.width = Pt(width)
    shape.shadow.inherit = False
    return shape


def add_picture_cover(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = w / h
    if image_ratio > frame_ratio:
        shown_w = h * image_ratio
        crop = (shown_w - w) / shown_w / 2
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        shown_h = w / image_ratio
        crop = (shown_h - h) / shown_h / 2
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
        picture.crop_top = crop
        picture.crop_bottom = crop
    picture.left = Inches(x)
    picture.top = Inches(y)
    picture.width = Inches(w)
    picture.height = Inches(h)
    return picture


def make_pptx() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(COLORS["bg"])

    # Header
    add_box(slide, 0.45, 0.29, 0.29, 0.29, fill="indigo", line="indigo", radius=True)
    add_box(slide, 0.535, 0.375, 0.12, 0.12, fill="surface", line="surface", radius=True)
    add_text(slide, "RoboBrain LLM：自然语言任务的可执行分解", 0.86, 0.20, 8.6, 0.47, size=24, bold=True)
    add_text(
        slide,
        "RoboAssembly  ·  UR5e assembly Template  →  fabrica_plumbers_block_ur5e_assembly",
        0.87,
        0.69,
        8.6,
        0.28,
        size=9.5,
        color="muted",
    )
    add_box(slide, 10.82, 0.30, 1.98, 0.40, fill="teal_light", line="teal_light")
    add_text(slide, "●  ONLINE LLM PLANNER", 10.92, 0.30, 1.78, 0.40, size=9, color="teal", bold=True, align=PP_ALIGN.CENTER)

    # Product UI panel
    add_text(slide, "对话式任务生成工作台", 0.48, 1.10, 3.2, 0.28, size=11, bold=True)
    add_text(slide, "资产库检索、规划过程与执行入口同屏呈现", 2.70, 1.10, 2.65, 0.28, size=8.5, color="muted", align=PP_ALIGN.RIGHT)
    add_box(slide, 0.45, 1.42, 5.02, 2.83, fill="surface", line="line", radius=True)
    add_picture_cover(slide, UI_FRAME, 0.50, 1.47, 4.92, 2.73)

    # Output evidence panel
    add_box(slide, 0.45, 4.48, 5.02, 2.55, fill="surface", line="line", radius=True)
    add_text(slide, "执行输出", 0.68, 4.67, 1.3, 0.28, size=11, bold=True)
    add_box(slide, 0.68, 5.10, 2.06, 1.42, fill="ink", line="ink", radius=True)
    add_picture_cover(slide, SIM_FRAME, 0.71, 5.13, 2.00, 1.36)
    add_box(slide, 0.80, 6.29, 1.65, 0.27, fill="ink", line="ink")
    add_text(slide, "Isaac Sim 双臂轨迹回放", 0.82, 6.29, 1.61, 0.27, size=7.7, color="surface", bold=True, align=PP_ALIGN.CENTER)

    metrics = [("15", "物体"), ("22", "目标"), ("36", "执行单元")]
    for index, (value, label) in enumerate(metrics):
        x = 2.96 + index * 0.74
        add_text(slide, value, x, 5.08, 0.62, 0.37, size=19, color="indigo", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 5.44, 0.62, 0.24, size=8, color="muted", align=PP_ALIGN.CENTER)
    add_text(slide, "每条轨迹步骤均包含", 2.96, 5.85, 2.18, 0.25, size=8.5, color="muted")
    add_text(slide, "Skill  ·  操作臂  ·  对象  ·  目标位姿", 2.96, 6.10, 2.20, 0.25, size=9.2, bold=True)
    add_text(slide, "夹爪状态  ·  完成条件  ·  静止校验", 2.96, 6.38, 2.20, 0.25, size=9.2, bold=True)

    # Prompt card
    right_x = 5.78
    right_w = 7.10
    add_box(slide, right_x, 1.10, right_w, 0.68, fill="surface", line="line", radius=True)
    add_box(slide, right_x + 0.15, 1.26, 0.56, 0.28, fill="indigo_light", line="indigo_light")
    add_text(slide, "USER", right_x + 0.15, 1.26, 0.56, 0.28, size=8.3, color="indigo", bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "使用 UR5e 双臂完成 plumbers-block 装配，并生成可执行轨迹",
        right_x + 0.84,
        1.18,
        6.02,
        0.48,
        size=12.5,
        bold=True,
    )

    # LLM pipeline
    pipeline = [
        ("01", "理解任务"),
        ("02", "匹配资产"),
        ("03", "生成 Menu\n与 Annotation"),
        ("04", "拆分子任务"),
        ("05", "编排 Skill"),
        ("06", "校验并执行"),
    ]
    step_w = 1.03
    gap = 0.16
    for index, (number, label) in enumerate(pipeline):
        x = right_x + index * (step_w + gap)
        fill = "teal_light" if index < 3 else "indigo_light"
        accent = "teal" if index < 3 else "indigo"
        add_box(slide, x, 1.95, step_w, 0.70, fill=fill, line=fill, radius=True)
        add_text(slide, number, x + 0.08, 2.03, 0.25, 0.22, size=8, color=accent, bold=True)
        add_text(slide, label, x + 0.08, 2.24, 0.87, 0.31, size=8.6, color="ink", bold=True, valign=MSO_ANCHOR.TOP)
        if index < len(pipeline) - 1:
            add_text(slide, "›", x + step_w, 2.13, gap, 0.30, size=17, color="muted", bold=True, align=PP_ALIGN.CENTER)

    # Grounded decomposition
    add_text(slide, "LLM 生成的真实任务分解", right_x, 2.88, 3.2, 0.30, size=11.5, bold=True)
    add_text(slide, "顺序约束：先稳定 block_2，再插装与堆叠", 9.46, 2.88, 3.40, 0.30, size=8.5, color="muted", align=PP_ALIGN.RIGHT)
    rows = [
        ("1", "建立稳定基座", "右臂", "block_2", "接近 → 抓取 → 抬升 → 转运 → 释放/锁定", "8"),
        ("2", "插入主槽零件", "左臂", "block_0", "接近 → 抓取 → 对位 → 插入 → 稳定", "7"),
        ("3", "完成上层堆叠", "左臂", "block_3", "接近 → 抓取 → 抬升 → 堆叠 → 稳定", "7"),
        ("4", "插入左孔零件", "左臂", "block_4", "接近 → 抓取 → 抬升 → 插入 → 稳定", "7"),
        ("5", "插入右孔并验收", "左臂", "block_1", "接近 → 抓取 → 抬升 → 插入 → 静止验收", "7"),
    ]
    row_y = 3.28
    row_h = 0.59
    row_gap = 0.09
    for index, (number, title, arm, obj, skills, count) in enumerate(rows):
        y = row_y + index * (row_h + row_gap)
        add_box(slide, right_x, y, right_w, row_h, fill="surface", line="line", radius=True)
        accent = "teal" if arm == "右臂" else "indigo"
        light = "teal_light" if arm == "右臂" else "indigo_light"
        add_box(slide, right_x + 0.14, y + 0.14, 0.30, 0.30, fill=accent, line=accent, radius=True)
        add_text(slide, number, right_x + 0.14, y + 0.14, 0.30, 0.30, size=8.5, color="surface", bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, right_x + 0.56, y + 0.08, 1.46, 0.43, size=10.2, bold=True)
        add_box(slide, right_x + 2.12, y + 0.16, 0.55, 0.27, fill=light, line=light)
        add_text(slide, arm, right_x + 2.12, y + 0.16, 0.55, 0.27, size=8.0, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, obj, right_x + 2.82, y + 0.08, 0.82, 0.43, size=9.1, color="muted", bold=True)
        add_text(slide, skills, right_x + 3.67, y + 0.08, 2.90, 0.43, size=8.8)
        add_box(slide, right_x + 6.62, y + 0.14, 0.32, 0.30, fill="green_light", line="green_light")
        add_text(slide, count, right_x + 6.62, y + 0.14, 0.32, 0.30, size=8.5, color="green", bold=True, align=PP_ALIGN.CENTER)

    # Footer callout
    add_box(slide, right_x, 6.75, right_w, 0.28, fill="ink", line="ink", radius=True)
    add_text(
        slide,
        "LLM 负责语义分解与资源映射  →  RoboChecker 校验约束  →  Isaac Sim 逐 Skill 执行",
        right_x + 0.18,
        6.75,
        right_w - 0.36,
        0.28,
        size=8.2,
        color="surface",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 2: next-month roadmap
    plan_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    plan_background = plan_slide.background.fill
    plan_background.solid()
    plan_background.fore_color.rgb = rgb(COLORS["bg"])

    add_box(plan_slide, 0.45, 0.29, 0.29, 0.29, fill="indigo", line="indigo", radius=True)
    add_box(plan_slide, 0.535, 0.375, 0.12, 0.12, fill="surface", line="surface", radius=True)
    add_text(plan_slide, "下月计划：从演示流程走向稳定自动化", 0.86, 0.20, 8.6, 0.47, size=24, bold=True)
    add_text(
        plan_slide,
        "核心目标：自然语言任务  →  自动生成新任务  →  仿真执行  →  可复现实验",
        0.87,
        0.69,
        8.8,
        0.28,
        size=9.5,
        color="muted",
    )
    add_box(plan_slide, 11.18, 0.30, 1.62, 0.40, fill="indigo_light", line="indigo_light")
    add_text(plan_slide, "NEXT 30 DAYS", 11.28, 0.30, 1.42, 0.40, size=9, color="indigo", bold=True, align=PP_ALIGN.CENTER)

    add_box(plan_slide, 0.45, 1.10, 12.43, 0.78, fill="surface", line="line", radius=True)
    add_box(plan_slide, 0.67, 1.31, 0.38, 0.36, fill="teal", line="teal", radius=True)
    add_text(plan_slide, "✓", 0.67, 1.31, 0.38, 0.36, size=12, color="surface", bold=True, align=PP_ALIGN.CENTER)
    add_text(
        plan_slide,
        "30 天目标：让 RoboBrain LLM 不只生成一次演示，而是稳定生成、校验并执行多类双臂装配任务",
        1.19,
        1.22,
        7.25,
        0.47,
        size=12.5,
        bold=True,
    )
    for index, (label, value) in enumerate([("任务扩展", "≥ 3 类"), ("批量运行", "≥ 10 条"), ("生成成功率", "≥ 80%")]):
        x = 8.68 + index * 1.31
        add_text(plan_slide, value, x, 1.17, 1.10, 0.30, size=12, color="indigo", bold=True, align=PP_ALIGN.CENTER)
        add_text(plan_slide, label, x, 1.46, 1.10, 0.22, size=8, color="muted", align=PP_ALIGN.CENTER)

    add_text(plan_slide, "四周推进路线", 0.47, 2.12, 2.2, 0.31, size=12, bold=True)
    add_text(plan_slide, "每周形成可演示、可检查的增量交付", 9.75, 2.12, 3.10, 0.31, size=8.8, color="muted", align=PP_ALIGN.RIGHT)

    weeks = [
        {
            "week": "WEEK 1",
            "title": "Planner 输出稳定化",
            "accent": "teal",
            "light": "teal_light",
            "items": ["接入真实 LLM 请求与状态流", "固定 Menu / Annotation / Skill Schema", "补齐超时、重试与错误提示"],
            "deliverable": "交付：稳定的在线规划接口",
        },
        {
            "week": "WEEK 2",
            "title": "资产检索与任务扩展",
            "accent": "indigo",
            "light": "indigo_light",
            "items": ["建立机器人、场景、物体语义索引", "自动生成任务目录与标注文件", "新增插装、堆叠、对接任务"],
            "deliverable": "交付：≥ 3 类新任务",
        },
        {
            "week": "WEEK 3",
            "title": "仿真闭环与恢复",
            "accent": "amber",
            "light": "amber_light",
            "items": ["逐 Skill 启动 Isaac Sim 执行", "实时回传阶段、对象与机械臂状态", "失败定位、重试与安全停止"],
            "deliverable": "交付：可观测执行闭环",
        },
        {
            "week": "WEEK 4",
            "title": "批量评测与交付",
            "accent": "green",
            "light": "green_light",
            "items": ["批量生成并回放 10 条以上轨迹", "统计生成、校验与执行成功率", "完善 MP4 录制和实验归档"],
            "deliverable": "交付：月度 Demo + 指标报告",
        },
    ]
    card_w = 2.93
    card_gap = 0.18
    card_y = 2.56
    for index, week in enumerate(weeks):
        x = 0.45 + index * (card_w + card_gap)
        add_box(plan_slide, x, card_y, card_w, 3.32, fill="surface", line="line", radius=True)
        add_box(plan_slide, x, card_y, card_w, 0.10, fill=week["accent"], line=week["accent"], radius=False)
        add_box(plan_slide, x + 0.20, card_y + 0.27, 0.72, 0.29, fill=week["light"], line=week["light"])
        add_text(plan_slide, week["week"], x + 0.20, card_y + 0.27, 0.72, 0.29, size=8.2, color=week["accent"], bold=True, align=PP_ALIGN.CENTER)
        add_text(plan_slide, week["title"], x + 0.20, card_y + 0.70, card_w - 0.40, 0.38, size=13, bold=True)
        for item_index, item in enumerate(week["items"]):
            item_y = card_y + 1.27 + item_index * 0.52
            add_box(plan_slide, x + 0.22, item_y + 0.08, 0.18, 0.18, fill=week["light"], line=week["light"], radius=True)
            add_text(plan_slide, str(item_index + 1), x + 0.22, item_y + 0.08, 0.18, 0.18, size=6.7, color=week["accent"], bold=True, align=PP_ALIGN.CENTER)
            add_text(plan_slide, item, x + 0.50, item_y, card_w - 0.72, 0.36, size=9.0)
        add_box(plan_slide, x + 0.20, card_y + 2.78, card_w - 0.40, 0.34, fill=week["light"], line=week["light"])
        add_text(
            plan_slide,
            week["deliverable"],
            x + 0.29,
            card_y + 2.78,
            card_w - 0.58,
            0.34,
            size=8.3,
            color=week["accent"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    add_text(plan_slide, "月末验收", 0.47, 6.13, 1.3, 0.28, size=11.5, bold=True)
    acceptance = [
        ("01", "自然语言输入可稳定生成完整任务文件"),
        ("02", "RoboChecker 可定位资源与约束错误"),
        ("03", "界面实时展示仿真阶段和失败原因"),
        ("04", "轨迹、日志与 MP4 可一键归档"),
    ]
    for index, (number, label) in enumerate(acceptance):
        x = 1.72 + index * 2.77
        add_text(plan_slide, number, x, 6.14, 0.34, 0.28, size=8.5, color="indigo", bold=True)
        add_text(plan_slide, label, x + 0.39, 6.10, 2.25, 0.38, size=8.5, bold=True)
    add_box(plan_slide, 0.45, 6.72, 12.43, 0.31, fill="ink", line="ink", radius=True)
    add_text(
        plan_slide,
        "关键路径：Planner Schema 稳定  →  资产语义对齐  →  Skill 执行可观测  →  批量评测闭环",
        0.70,
        6.72,
        11.93,
        0.31,
        size=8.6,
        color="surface",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    presentation.save(PPTX_PATH)


def pil_font(size: int, bold: bool = False):
    path = FONT_BOLD_PATH if bold and Path(FONT_BOLD_PATH).exists() else FONT_PATH
    return ImageFont.truetype(path, size=size, index=2)


def make_preview() -> None:
    """Create a faithful raster preview for quick visual inspection."""
    scale_x = PX_W / SLIDE_W
    scale_y = PX_H / SLIDE_H

    def box(coords, fill, outline=None, radius=16, width=2):
        x, y, w, h = coords
        xy = (int(x * scale_x), int(y * scale_y), int((x + w) * scale_x), int((y + h) * scale_y))
        draw.rounded_rectangle(xy, radius=radius, fill="#" + COLORS.get(fill, fill), outline=("#" + COLORS.get(outline, outline)) if outline else None, width=width)

    def text_at(value, x, y, size, fill="ink", bold=False, anchor="la"):
        draw.text((int(x * scale_x), int(y * scale_y)), value, font=pil_font(size, bold), fill="#" + COLORS.get(fill, fill), anchor=anchor, spacing=3)

    def paste_cover(path: Path, x, y, w, h):
        image = Image.open(path).convert("RGB")
        target_w, target_h = int(w * scale_x), int(h * scale_y)
        ratio = max(target_w / image.width, target_h / image.height)
        resized = image.resize((int(image.width * ratio), int(image.height * ratio)), Image.Resampling.LANCZOS)
        left = (resized.width - target_w) // 2
        top = (resized.height - target_h) // 2
        cropped = resized.crop((left, top, left + target_w, top + target_h))
        canvas.paste(cropped, (int(x * scale_x), int(y * scale_y)))

    canvas = Image.new("RGB", (PX_W, PX_H), "#" + COLORS["bg"])
    draw = ImageDraw.Draw(canvas)

    box((0.45, 0.29, 0.29, 0.29), "indigo", "indigo", 12)
    box((0.535, 0.375, 0.12, 0.12), "surface", "surface", 6)
    text_at("RoboBrain LLM：自然语言任务的可执行分解", 0.86, 0.49, 46, "ink", True, "lm")
    text_at("RoboAssembly  ·  UR5e assembly Template  →  fabrica_plumbers_block_ur5e_assembly", 0.87, 0.82, 18, "muted", False, "lm")
    box((10.82, 0.30, 1.98, 0.40), "teal_light", "teal_light", 18)
    text_at("●  ONLINE LLM PLANNER", 11.81, 0.50, 16, "teal", True, "mm")

    text_at("对话式任务生成工作台", 0.48, 1.24, 21, "ink", True, "lm")
    text_at("资产库检索、规划过程与执行入口同屏呈现", 5.34, 1.24, 16, "muted", False, "rm")
    box((0.45, 1.42, 5.02, 2.83), "surface", "line", 18)
    paste_cover(UI_FRAME, 0.50, 1.47, 4.92, 2.73)

    box((0.45, 4.48, 5.02, 2.55), "surface", "line", 18)
    text_at("执行输出", 0.68, 4.81, 21, "ink", True, "lm")
    box((0.68, 5.10, 2.06, 1.42), "ink", "ink", 16)
    paste_cover(SIM_FRAME, 0.71, 5.13, 2.00, 1.36)
    box((0.80, 6.29, 1.65, 0.27), "ink", "ink", 8)
    text_at("Isaac Sim 双臂轨迹回放", 1.625, 6.425, 14, "surface", True, "mm")
    for index, (value, label) in enumerate([("15", "物体"), ("22", "目标"), ("36", "执行单元")]):
        x = 3.27 + index * 0.74
        text_at(value, x, 5.33, 37, "indigo", True, "mm")
        text_at(label, x, 5.59, 14, "muted", False, "mm")
    text_at("每条轨迹步骤均包含", 2.96, 5.98, 16, "muted", False, "lm")
    text_at("Skill  ·  操作臂  ·  对象  ·  目标位姿", 2.96, 6.23, 17, "ink", True, "lm")
    text_at("夹爪状态  ·  完成条件  ·  静止校验", 2.96, 6.51, 17, "ink", True, "lm")

    right_x = 5.78
    right_w = 7.10
    box((right_x, 1.10, right_w, 0.68), "surface", "line", 18)
    box((right_x + 0.15, 1.26, 0.56, 0.28), "indigo_light", "indigo_light", 10)
    text_at("USER", right_x + 0.43, 1.40, 15, "indigo", True, "mm")
    text_at("使用 UR5e 双臂完成 plumbers-block 装配，并生成可执行轨迹", right_x + 0.84, 1.43, 24, "ink", True, "lm")

    pipeline = [("01", "理解任务"), ("02", "匹配资产"), ("03", "生成 Menu\n与 Annotation"), ("04", "拆分子任务"), ("05", "编排 Skill"), ("06", "校验并执行")]
    step_w, gap = 1.03, 0.16
    for index, (number, label) in enumerate(pipeline):
        x = right_x + index * (step_w + gap)
        fill = "teal_light" if index < 3 else "indigo_light"
        accent = "teal" if index < 3 else "indigo"
        box((x, 1.95, step_w, 0.70), fill, fill, 14)
        text_at(number, x + 0.08, 2.11, 14, accent, True, "lm")
        text_at(label, x + 0.08, 2.31, 16, "ink", True, "la")
        if index < len(pipeline) - 1:
            text_at("›", x + step_w + gap / 2, 2.30, 29, "muted", True, "mm")

    text_at("LLM 生成的真实任务分解", right_x, 3.03, 22, "ink", True, "lm")
    text_at("顺序约束：先稳定 block_2，再插装与堆叠", 12.86, 3.03, 15, "muted", False, "rm")
    rows = [
        ("1", "建立稳定基座", "右臂", "block_2", "接近 → 抓取 → 抬升 → 转运 → 释放/锁定", "8"),
        ("2", "插入主槽零件", "左臂", "block_0", "接近 → 抓取 → 对位 → 插入 → 稳定", "7"),
        ("3", "完成上层堆叠", "左臂", "block_3", "接近 → 抓取 → 抬升 → 堆叠 → 稳定", "7"),
        ("4", "插入左孔零件", "左臂", "block_4", "接近 → 抓取 → 抬升 → 插入 → 稳定", "7"),
        ("5", "插入右孔并验收", "左臂", "block_1", "接近 → 抓取 → 抬升 → 插入 → 静止验收", "7"),
    ]
    for index, (number, title, arm, obj, skills, count) in enumerate(rows):
        y = 3.28 + index * 0.68
        box((right_x, y, right_w, 0.59), "surface", "line", 14)
        accent = "teal" if arm == "右臂" else "indigo"
        light = "teal_light" if arm == "右臂" else "indigo_light"
        box((right_x + 0.14, y + 0.14, 0.30, 0.30), accent, accent, 14)
        text_at(number, right_x + 0.29, y + 0.29, 16, "surface", True, "mm")
        text_at(title, right_x + 0.56, y + 0.30, 19, "ink", True, "lm")
        box((right_x + 2.12, y + 0.16, 0.55, 0.27), light, light, 10)
        text_at(arm, right_x + 2.395, y + 0.295, 14, accent, True, "mm")
        text_at(obj, right_x + 2.82, y + 0.30, 16, "muted", True, "lm")
        text_at(skills, right_x + 3.67, y + 0.30, 16, "ink", False, "lm")
        box((right_x + 6.62, y + 0.14, 0.32, 0.30), "green_light", "green_light", 10)
        text_at(count, right_x + 6.78, y + 0.29, 16, "green", True, "mm")

    box((right_x, 6.75, right_w, 0.28), "ink", "ink", 10)
    text_at("LLM 负责语义分解与资源映射  →  RoboChecker 校验约束  →  Isaac Sim 逐 Skill 执行", right_x + right_w / 2, 6.89, 15, "surface", True, "mm")
    canvas.save(PREVIEW_PATH, quality=95)


def make_plan_preview() -> None:
    """Create a raster preview of the next-month roadmap slide."""
    scale_x = PX_W / SLIDE_W
    scale_y = PX_H / SLIDE_H
    canvas = Image.new("RGB", (PX_W, PX_H), "#" + COLORS["bg"])
    draw = ImageDraw.Draw(canvas)

    def box(coords, fill, outline=None, radius=16, width=2):
        x, y, w, h = coords
        xy = (int(x * scale_x), int(y * scale_y), int((x + w) * scale_x), int((y + h) * scale_y))
        draw.rounded_rectangle(
            xy,
            radius=radius,
            fill="#" + COLORS.get(fill, fill),
            outline=("#" + COLORS.get(outline, outline)) if outline else None,
            width=width,
        )

    def text_at(value, x, y, size, fill="ink", bold=False, anchor="la"):
        draw.text((int(x * scale_x), int(y * scale_y)), value, font=pil_font(size, bold), fill="#" + COLORS.get(fill, fill), anchor=anchor, spacing=3)

    box((0.45, 0.29, 0.29, 0.29), "indigo", "indigo", 12)
    box((0.535, 0.375, 0.12, 0.12), "surface", "surface", 6)
    text_at("下月计划：从演示流程走向稳定自动化", 0.86, 0.49, 46, "ink", True, "lm")
    text_at("核心目标：自然语言任务  →  自动生成新任务  →  仿真执行  →  可复现实验", 0.87, 0.82, 18, "muted", False, "lm")
    box((11.18, 0.30, 1.62, 0.40), "indigo_light", "indigo_light", 18)
    text_at("NEXT 30 DAYS", 11.99, 0.50, 16, "indigo", True, "mm")

    box((0.45, 1.10, 12.43, 0.78), "surface", "line", 18)
    box((0.67, 1.31, 0.38, 0.36), "teal", "teal", 14)
    text_at("✓", 0.86, 1.49, 22, "surface", True, "mm")
    text_at("30 天目标：让 RoboBrain LLM 不只生成一次演示，而是稳定生成、校验并执行多类双臂装配任务", 1.19, 1.47, 24, "ink", True, "lm")
    for index, (label, value) in enumerate([("任务扩展", "≥ 3 类"), ("批量运行", "≥ 10 条"), ("生成成功率", "≥ 80%")]):
        x = 9.23 + index * 1.31
        text_at(value, x, 1.35, 23, "indigo", True, "mm")
        text_at(label, x, 1.62, 14, "muted", False, "mm")

    text_at("四周推进路线", 0.47, 2.27, 23, "ink", True, "lm")
    text_at("每周形成可演示、可检查的增量交付", 12.85, 2.27, 16, "muted", False, "rm")
    weeks = [
        ("WEEK 1", "Planner 输出稳定化", "teal", "teal_light", ["接入真实 LLM 请求与状态流", "固定 Menu / Annotation / Skill Schema", "补齐超时、重试与错误提示"], "交付：稳定的在线规划接口"),
        ("WEEK 2", "资产检索与任务扩展", "indigo", "indigo_light", ["建立机器人、场景、物体语义索引", "自动生成任务目录与标注文件", "新增插装、堆叠、对接任务"], "交付：≥ 3 类新任务"),
        ("WEEK 3", "仿真闭环与恢复", "amber", "amber_light", ["逐 Skill 启动 Isaac Sim 执行", "实时回传阶段、对象与机械臂状态", "失败定位、重试与安全停止"], "交付：可观测执行闭环"),
        ("WEEK 4", "批量评测与交付", "green", "green_light", ["批量生成并回放 10 条以上轨迹", "统计生成、校验与执行成功率", "完善 MP4 录制和实验归档"], "交付：月度 Demo + 指标报告"),
    ]
    card_w, card_gap, card_y = 2.93, 0.18, 2.56
    for index, (week, title, accent, light, items, deliverable) in enumerate(weeks):
        x = 0.45 + index * (card_w + card_gap)
        box((x, card_y, card_w, 3.32), "surface", "line", 18)
        draw.rectangle(
            (int(x * scale_x), int(card_y * scale_y), int((x + card_w) * scale_x), int((card_y + 0.10) * scale_y)),
            fill="#" + COLORS[accent],
        )
        box((x + 0.20, card_y + 0.27, 0.72, 0.29), light, light, 10)
        text_at(week, x + 0.56, card_y + 0.415, 15, accent, True, "mm")
        text_at(title, x + 0.20, card_y + 0.88, 25, "ink", True, "lm")
        for item_index, item in enumerate(items):
            item_y = card_y + 1.27 + item_index * 0.52
            box((x + 0.22, item_y + 0.08, 0.18, 0.18), light, light, 8)
            text_at(str(item_index + 1), x + 0.31, item_y + 0.17, 12, accent, True, "mm")
            text_at(item, x + 0.50, item_y + 0.18, 17, "ink", False, "lm")
        box((x + 0.20, card_y + 2.78, card_w - 0.40, 0.34), light, light, 10)
        text_at(deliverable, x + card_w / 2, card_y + 2.95, 15, accent, True, "mm")

    text_at("月末验收", 0.47, 6.27, 22, "ink", True, "lm")
    acceptance = [
        ("01", "自然语言输入可稳定生成完整任务文件"),
        ("02", "RoboChecker 可定位资源与约束错误"),
        ("03", "界面实时展示仿真阶段和失败原因"),
        ("04", "轨迹、日志与 MP4 可一键归档"),
    ]
    for index, (number, label) in enumerate(acceptance):
        x = 1.72 + index * 2.77
        text_at(number, x, 6.28, 16, "indigo", True, "lm")
        text_at(label, x + 0.39, 6.28, 16, "ink", True, "lm")
    box((0.45, 6.72, 12.43, 0.31), "ink", "ink", 10)
    text_at("关键路径：Planner Schema 稳定  →  资产语义对齐  →  Skill 执行可观测  →  批量评测闭环", 6.665, 6.875, 16, "surface", True, "mm")
    canvas.save(PLAN_PREVIEW_PATH, quality=95)


if __name__ == "__main__":
    if not UI_FRAME.exists() or not SIM_FRAME.exists():
        raise SystemExit("Presentation source frames are missing. Extract frame_035.png and frame_175.png first.")
    make_pptx()
    make_preview()
    make_plan_preview()
    print(PPTX_PATH)
    print(PREVIEW_PATH)
    print(PLAN_PREVIEW_PATH)
